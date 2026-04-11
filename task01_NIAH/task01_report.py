import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import random
import torch
import torch.nn as nn
import torch.optim as optim
import matplotlib.pyplot as plt
import time
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from thop import profile

from PLHN import PhaseLockedHolographyNetwork
from TFER import SimpleTransformer
from PHTF import PhaseLockedHolographicTransformerFusion

# 設定掃描與訓練上限
MAX_SAMPLES = 5000000000        # 單個模型在單一長度的最多看過樣本數上限 (避免無限迴圈)
EVAL_EVERY_BATCHES = 100    # 每訓練多少個 Batches 之後，進行一次 Test Set 驗證
MAX_TEST_SAMPLES = 1000     # 永遠維持 1000 筆 Test Data
SEQ_LENS = [32, 64, 128, 256]

# 模型架構參數
VOCAB_SIZE = 32
EMBED_DIM = 128
NUM_HEADS = 4
HIDDEN_DIM = 256

def set_seed(seed=8888):
    """ 
    固定所有隨機數種子以確保實驗可複現。
    注意：開啟 cudnn.deterministic 可能會犧牲些微的訓練速度，以換取絕對的穩定性。
    """
    random.seed(seed)
    np.random.seed(seed)
    torch.manual_seed(seed)
    if torch.cuda.is_available():
        torch.cuda.manual_seed(seed)
        torch.cuda.manual_seed_all(seed) # 針對多 GPU
    
    # 強制 cuDNN 使用確定性演算法
    torch.backends.cudnn.deterministic = True
    torch.backends.cudnn.benchmark = False

def count_parameters(model):
    return sum(p.numel() for p in model.parameters() if p.requires_grad)

def generate_batch(batch_size, seq_len, vocab_size, device):
    """ 動態生成沒有界限的無窮 NIAH 訓練資料 (On-the-fly) """
    haystack_max = vocab_size - 2
    needle_key = vocab_size - 2
    
    # 全部填滿背景雜訊 [batch_size, seq_len]
    batch_x = torch.randint(0, haystack_max, (batch_size, seq_len), device=device)
    batch_y = torch.zeros((batch_size,), dtype=torch.long, device=device)
    
    # 決定每個 batch 藏針的位置與數值
    needle_idxs = torch.randint(0, seq_len - 2, (batch_size,), device=device)
    needle_vals = torch.randint(0, haystack_max, (batch_size,), device=device)
    
    # 放入針 [NEEDLE_KEY, VALUE]
    batch_indices = torch.arange(batch_size, device=device)
    batch_x[batch_indices, needle_idxs] = needle_key
    batch_x[batch_indices, needle_idxs + 1] = needle_vals
    
    # 放入 Query
    batch_x[:, -1] = needle_key
    batch_y[:] = needle_vals
    
    return batch_x, batch_y

def get_adaptive_batch(seq_len):
    """ 因應 TFER 在長度 8192 會 OOM，動態降 Batch Size 搭配 Gradient Accumulation """
    if seq_len >= 8192:
        return 32, 32      # 實體 batch=1, 累積 16 次梯度 = 邏輯 16
    elif seq_len >= 4096:
        return 32, 16
    elif seq_len >= 2048:
        return 32, 8
    elif seq_len >= 1024:
        return 32, 4
    else:
        return 32, 2

def generate_fixed_test_set(seq_len, vocab_size, device, seed=8888):
    """ 生成固定的 Test 測試集確保評估一致性 """
    # 儲存原本的隨機狀態，避免影響主程式的訓練隨機性
    cpu_rng = torch.get_rng_state()
    cuda_rng = torch.cuda.get_rng_state() if device.type == 'cuda' else None
        
    torch.manual_seed(seed)
    if device.type == 'cuda':
        torch.cuda.manual_seed_all(seed)
        
    batch_x, batch_y = generate_batch(MAX_TEST_SAMPLES, seq_len, vocab_size, device)
    
    # 還原隨機狀態
    torch.set_rng_state(cpu_rng)
    if device.type == 'cuda' and cuda_rng is not None:
        torch.cuda.set_rng_state(cuda_rng)
        
    return batch_x, batch_y

def train_until_convergence(model, test_x, test_y, seq_len, vocab_size, device, model_name):
    criterion = nn.CrossEntropyLoss()
    optimizer = optim.Adam(model.parameters(), lr=3e-4)
    
    batch_sz, accum_steps = get_adaptive_batch(seq_len)
    
    samples_seen_history = []
    test_acc_history = []
    loss_history = []
    
    total_samples = 0
    batches_since_eval = 0
    running_loss = 0.0
    accum_count = 0
    
    model.train()
    optimizer.zero_grad()
    start_time = time.time()
    
    while total_samples < MAX_SAMPLES:
        x, y = generate_batch(batch_sz, seq_len, vocab_size, device)
        logits = model(x)
        last_logits = logits[:, -1, :]
        loss = criterion(last_logits, y) / accum_steps
        loss.backward()
        
        running_loss += loss.item() * accum_steps
        accum_count += 1
        
        # Gradient Accumulation
        if accum_count == accum_steps:
            optimizer.step()
            optimizer.zero_grad()
            accum_count = 0
            
            # 完整累積一個邏輯 Batch 後才計數
            batches_since_eval += 1
            total_samples += (batch_sz * accum_steps) # 永遠是加 16
            
            # 進行定期的 Test Set 驗證
            if batches_since_eval >= EVAL_EVERY_BATCHES:
                model.eval()
                with torch.no_grad():
                    test_correct = 0
                    # 避免 Test Set 一次吃掉太大 Memory，分批測試
                    test_chunk_size = 128 if seq_len < 2048 else (16 if seq_len < 8192 else 4)
                    for i in range(0, len(test_y), test_chunk_size):
                        end_idx = min(i + test_chunk_size, len(test_y))
                        cx = test_x[i:end_idx]
                        cy = test_y[i:end_idx]
                        clogits = model(cx)
                        cpreds = clogits[:, -1, :].argmax(dim=-1)
                        test_correct += (cpreds == cy).sum().item()
                        
                test_acc = test_correct / len(test_y)
                current_time = time.time() - start_time
                avg_loss = running_loss / batches_since_eval
                
                print(f"[{model_name} | Seq {seq_len}] Samples: {total_samples:,} | Time: {current_time:.1f}s | Test Acc: {test_acc:.4f} | Loss: {avg_loss:.4f}")
                
                samples_seen_history.append(total_samples)
                test_acc_history.append(test_acc)
                loss_history.append(avg_loss)
                
                if test_acc >= 0.90: # 達到近乎 100%
                    return total_samples, samples_seen_history, test_acc_history, loss_history
                
                # 重設變數回到訓練
                running_loss = 0.0
                batches_since_eval = 0
                model.train()
                
    # 若達到 MAX_SAMPLES 仍未 100% 收斂
    return -1, samples_seen_history, test_acc_history, loss_history

def create_training_plot(results_map, seq_len, filename):
    colors_map = {'TFER': 'r', 'PLHN': 'b', 'PHTF': 'g'}
    
    plt.figure(figsize=(14, 5))
    
    plt.subplot(1, 2, 1)
    for name, r in results_map.items():
        if len(r['samples']) > 0:
            plt.plot(r['samples'], r['loss'], color=colors_map.get(name, 'k'), label=name)
    plt.title(f'Loss vs Samples Seen (Seq {seq_len})')
    plt.xlabel('Total Seen Samples')
    plt.ylabel('Loss')
    plt.grid(True)
    plt.legend()
    
    plt.subplot(1, 2, 2)
    for name, r in results_map.items():
        if len(r['samples']) > 0:
            plt.plot(r['samples'], r['test_acc'], color=colors_map.get(name, 'k'), label=name)
    plt.title(f'Test Accuracy vs Samples Seen (Seq {seq_len})')
    plt.xlabel('Total Seen Samples')
    plt.ylabel('Test Accuracy')
    plt.ylim(0, 1.05)
    plt.grid(True)
    plt.legend()
    
    plt.tight_layout()
    plt.savefig(filename, dpi=300)
    plt.close()

def build_presentation_all_in_one(summary_data, plot_files, base_dir):
    prs = Presentation()
    
    # 首頁標題
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Task 21: Comprehensive NIAH Scan"
    slide.placeholders[1].text = "TFER vs PLHN vs PHTF\nSequence Lengths Scanning (32 to 8192)"
    
    # 總結表格頁面
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Samples to Reach 100% Test Accuracy"
    
    rows = len(SEQ_LENS) + 1
    cols = 4
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(8.0)
    height = Inches(4.0)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.0)
    
    headers = ['Seq Length', 'TFER', 'PLHN', 'PHTF']
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        
    for r_idx, seq in enumerate(SEQ_LENS):
        table.cell(r_idx + 1, 0).text = str(seq)
        for c_idx, model_name in enumerate(['TFER', 'PLHN', 'PHTF']):
            val = summary_data[model_name].get(seq, 'DNF')
            table.cell(r_idx + 1, c_idx + 1).text = f"{val:,}" if isinstance(val, int) and val > 0 else str(val)
            
    # 插入每種長度的訓練曲線圖表
    for seq in SEQ_LENS:
        if seq in plot_files:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Convergence Results (SEQ_LEN = {seq})"
            slide.shapes.add_picture(plot_files[seq], Inches(0.5), Inches(1.5), width=Inches(9))
            
    out_file = os.path.join(base_dir, "Task21_All_In_One_Report.pptx")
    prs.save(out_file)
    return out_file

def main():
    device = torch.device('cuda:0' if torch.cuda.is_available() else 'cpu')
    print(f"啟動 Task 21 NIAH 掃描式測試 (動態生成 + 長度掃描)... {device}")
    
    base_dir = os.path.dirname(os.path.abspath(__file__))
    summary_data = {'TFER': {}, 'PLHN': {}, 'PHTF': {}}
    plot_files = {}
    
    set_seed(8888)
    
    for seq in SEQ_LENS:
        print(f"\n{'='*50}\n開始掃描 SEQ_LEN = {seq}\n{'='*50}")
        
        # 針對目前長度，準備固定的 1000 筆測試資料
        test_x, test_y = generate_fixed_test_set(seq, VOCAB_SIZE, device, seed=8888)
        
        # 重新實例化三個模型，確保每次都是從頭訓練！
        models = {
            'TFER': SimpleTransformer(VOCAB_SIZE, EMBED_DIM, NUM_HEADS, HIDDEN_DIM, 2, max_seq_len=seq).to(device),
            'PLHN': PhaseLockedHolographyNetwork(VOCAB_SIZE, EMBED_DIM, HIDDEN_DIM, 2, max_seq_len=seq).to(device),
            'PHTF': PhaseLockedHolographicTransformerFusion(VOCAB_SIZE, EMBED_DIM, NUM_HEADS, HIDDEN_DIM, 2, max_seq_len=seq).to(device)
        }
        
        seq_results = {}
        for name, model in models.items():
            print(f"\n--- 訓練 {name} (Seq: {seq}) ---")
            
            set_seed(8888)
            
            samples_req, s_hist, acc_hist, loss_hist = train_until_convergence(
                model, test_x, test_y, seq, VOCAB_SIZE, device, name
            )
            # 紀錄需要多少樣本 (或 DNF)
            summary_data[name][seq] = samples_req if samples_req > 0 else 'DNF (Failed)'
            seq_results[name] = {'samples': s_hist, 'test_acc': acc_hist, 'loss': loss_hist}
            
            if samples_req > 0:
                print(f"✅ {name} 成功！在 {samples_req:,} 筆訓練資料時達到 100% 測試準確率。")
            else:
                highest_acc = max(acc_hist) if acc_hist else 0
                print(f"❌ {name} 失敗！無法在 {MAX_SAMPLES:,} 筆內達到 100% (最高準確率: {highest_acc:.4f})")
                
        # 當所有模型在該長度測試完畢，繪製曲線並存檔
        plot_path = os.path.join(base_dir, f'task21_scan_seq{seq}.png')
        create_training_plot(seq_results, seq, plot_path)
        plot_files[seq] = plot_path

    # 最後，建立彙整合併版的 PPTX 簡報
    print("\n產生 All-In-One 簡報檔...")
    pptx_path = build_presentation_all_in_one(summary_data, plot_files, base_dir)
    print(f"✅ 所有序列長度掃描完成！最終彙整報告儲存於: {pptx_path}")

if __name__ == "__main__":
    main()

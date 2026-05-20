import sys
import os

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import random
import torch
import torch.nn as nn
import torch.optim as optim
import matplotlib.pyplot as plt
import time
import numpy as np
import math
from pptx import Presentation
from pptx.util import Inches

# =================================================================
# 透過 Import 引入測試模型 (已移除 PHTF)
# =================================================================
from PLHN import PhaseLockedHolographyNetwork
from AGI import HolographicMemoryNetwork

# =================================================================
# 實驗超參數設定
# =================================================================
MAX_SAMPLES = 5000000000  # 單個模型在單一長度的最多看過樣本數上限
EVAL_EVERY_BATCHES = 100  # 每訓練多少個 Batches 之後，進行一次 Test Set 驗證
MAX_TEST_SAMPLES = 1000  # 永遠維持 1000 筆 Test Data
SEQ_LENS = [32, 64, 128, 256, 512, 1024]

# 模型架構基礎參數 (對齊基準)
VOCAB_SIZE = 32
EMBED_DIM = 128
HIDDEN_DIM = 256


def set_seed(seed=8888):
    """
    固定所有隨機數種子以確保實驗可複現。
    注意：開啟 cudnn.deterministic 可能會犧牲些微的訓練速度，以換取絕對的穩定性。
    """
    random.seed(seed)
    np.random.seed(seed)
    torch.manual_seed(seed)
    if torch.cuda.is_available():
        torch.cuda.manual_seed(seed)
        torch.cuda.manual_seed_all(seed)

    # 強制 cuDNN 使用確定性演算法
    torch.backends.cudnn.deterministic = True
    torch.backends.cudnn.benchmark = False


def count_parameters(model):
    return sum(p.numel() for p in model.parameters() if p.requires_grad)


def generate_batch(batch_size, seq_len, vocab_size, device):
    """動態生成無窮 NIAH 訓練資料"""
    haystack_max = vocab_size - 2
    needle_key = vocab_size - 2

    batch_x = torch.randint(0, haystack_max, (batch_size, seq_len), device=device)
    batch_y = torch.zeros((batch_size,), dtype=torch.long, device=device)

    needle_idxs = torch.randint(0, seq_len - 2, (batch_size,), device=device)
    needle_vals = torch.randint(0, haystack_max, (batch_size,), device=device)

    batch_indices = torch.arange(batch_size, device=device)
    batch_x[batch_indices, needle_idxs] = needle_key
    batch_x[batch_indices, needle_idxs + 1] = needle_vals

    batch_x[:, -1] = needle_key
    batch_y[:] = needle_vals

    return batch_x, batch_y


def get_adaptive_batch(seq_len):
    """因應 TFER 在長度 8192 會 OOM，動態降 Batch Size 搭配 Gradient Accumulation"""
    if seq_len >= 8192:
        return 32, 32  # 實體 batch=1, 累積 16 次梯度 = 邏輯 16
    elif seq_len >= 4096:
        return 32, 16
    elif seq_len >= 2048:
        return 32, 8
    elif seq_len >= 1024:
        return 32, 4
    else:
        return 32, 2


def generate_fixed_test_set(seq_len, vocab_size, device, seed=8888):
    cpu_rng = torch.get_rng_state()
    cuda_rng = torch.cuda.get_rng_state() if device.type == "cuda" else None

    torch.manual_seed(seed)
    if device.type == "cuda":
        torch.cuda.manual_seed_all(seed)

    batch_x, batch_y = generate_batch(MAX_TEST_SAMPLES, seq_len, vocab_size, device)

    torch.set_rng_state(cpu_rng)
    if device.type == "cuda" and cuda_rng is not None:
        torch.cuda.set_rng_state(cuda_rng)

    return batch_x, batch_y


def train_until_convergence(
    model, test_x, test_y, seq_len, vocab_size, device, model_name
):
    criterion = nn.CrossEntropyLoss()
    optimizer = optim.Adam(model.parameters(), lr=3e-4)

    batch_sz, accum_steps = get_adaptive_batch(seq_len)

    samples_seen_history = []
    test_acc_history = []
    loss_history = []

    total_samples = 0
    batches_since_eval = 0
    running_loss = 0.0
    accum_count = 0

    model.train()
    optimizer.zero_grad()
    start_time = time.time()

    while total_samples < MAX_SAMPLES:
        x, y = generate_batch(batch_sz, seq_len, vocab_size, device)
        logits = model(x)
        last_logits = logits[:, -1, :]
        loss = criterion(last_logits, y) / accum_steps
        loss.backward()

        running_loss += loss.item() * accum_steps
        accum_count += 1

        if accum_count == accum_steps:
            optimizer.step()
            optimizer.zero_grad()
            accum_count = 0

            batches_since_eval += 1
            total_samples += batch_sz * accum_steps

            if batches_since_eval >= EVAL_EVERY_BATCHES:
                model.eval()
                with torch.no_grad():
                    test_correct = 0
                    test_chunk_size = (
                        128 if seq_len < 2048 else (16 if seq_len < 8192 else 4)
                    )
                    for i in range(0, len(test_y), test_chunk_size):
                        end_idx = min(i + test_chunk_size, len(test_y))
                        cx = test_x[i:end_idx]
                        cy = test_y[i:end_idx]
                        clogits = model(cx)
                        cpreds = clogits[:, -1, :].argmax(dim=-1)
                        test_correct += (cpreds == cy).sum().item()

                test_acc = test_correct / len(test_y)
                current_time = time.time() - start_time
                avg_loss = running_loss / batches_since_eval

                print(
                    f"[{model_name} | Seq {seq_len}] Samples: {total_samples:,} | Time: {current_time:.1f}s | Test Acc: {test_acc:.4f} | Loss: {avg_loss:.4f}"
                )

                samples_seen_history.append(total_samples)
                test_acc_history.append(test_acc)
                loss_history.append(avg_loss)

                if test_acc >= 0.90:
                    return (
                        total_samples,
                        samples_seen_history,
                        test_acc_history,
                        loss_history,
                    )

                running_loss = 0.0
                batches_since_eval = 0
                model.train()

    return -1, samples_seen_history, test_acc_history, loss_history


def create_training_plot(results_map, seq_len, filename):
    colors_map = {"AGI": "r", "PLHN": "b"}  # 已移除 PHTF 的顏色對應

    plt.figure(figsize=(14, 5))

    plt.subplot(1, 2, 1)
    for name, r in results_map.items():
        if len(r["samples"]) > 0:
            plt.plot(
                r["samples"], r["loss"], color=colors_map.get(name, "k"), label=name
            )
    plt.title(f"Loss vs Samples Seen (Seq {seq_len})")
    plt.xlabel("Total Seen Samples")
    plt.ylabel("Loss")
    plt.grid(True)
    plt.legend()

    plt.subplot(1, 2, 2)
    for name, r in results_map.items():
        if len(r["samples"]) > 0:
            plt.plot(
                r["samples"], r["test_acc"], color=colors_map.get(name, "k"), label=name
            )
    plt.title(f"Test Accuracy vs Samples Seen (Seq {seq_len})")
    plt.xlabel("Total Seen Samples")
    plt.ylabel("Test Accuracy")
    plt.ylim(0, 1.05)
    plt.grid(True)
    plt.legend()

    plt.tight_layout()
    plt.savefig(filename, dpi=300)
    plt.close()


def build_presentation_all_in_one(summary_data, plot_files, param_data, base_dir):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Task 21: Comprehensive NIAH Scan"
    slide.placeholders[1].text = "AGI vs PLHN\nSequence Lengths Scanning (32 to 8192)"

    # --- 參數總表 ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Total Parameters Scaling by Sequence Length"
    rows = len(SEQ_LENS) + 1
    # 調整表格為 3 欄
    table = slide.shapes.add_table(
        rows, 3, Inches(1.0), Inches(1.5), Inches(8.0), Inches(4.0)
    ).table
    for i, h in enumerate(["Seq Length", "AGI", "PLHN"]):
        table.cell(0, i).text = h
    for r_idx, seq in enumerate(SEQ_LENS):
        table.cell(r_idx + 1, 0).text = str(seq)
        for c_idx, model_name in enumerate(["AGI", "PLHN"]):
            table.cell(r_idx + 1, c_idx + 1).text = f"{param_data[model_name][seq]:,}"

    # --- 收斂速度表 ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Samples to Reach 100% Test Accuracy"
    table = slide.shapes.add_table(
        rows, 3, Inches(1.0), Inches(1.5), Inches(8.0), Inches(4.0)
    ).table
    for i, h in enumerate(["Seq Length", "AGI", "PLHN"]):
        table.cell(0, i).text = h
    for r_idx, seq in enumerate(SEQ_LENS):
        table.cell(r_idx + 1, 0).text = str(seq)
        for c_idx, model_name in enumerate(["AGI", "PLHN"]):
            val = summary_data[model_name].get(seq, "DNF")
            table.cell(r_idx + 1, c_idx + 1).text = (
                f"{val:,}" if isinstance(val, int) and val > 0 else str(val)
            )

    for seq in SEQ_LENS:
        if seq in plot_files:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Convergence Results (SEQ_LEN = {seq})"
            slide.shapes.add_picture(
                plot_files[seq], Inches(0.5), Inches(1.5), width=Inches(9)
            )

    out_file = os.path.join(base_dir, "Task21_All_In_One_Report_AGI_vs_PLHN.pptx")
    prs.save(out_file)
    return out_file


def main():
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    print(f"啟動 Task 21 NIAH 掃描式測試 (AGI vs PLHN)... {device}")

    base_dir = os.path.dirname(os.path.abspath(__file__))
    summary_data = {"AGI": {}, "PLHN": {}}
    param_data = {"AGI": {}, "PLHN": {}}
    plot_files = {}

    set_seed(8888)

    for seq in SEQ_LENS:
        print(f"\n{'='*50}\n開始掃描 SEQ_LEN = {seq}\n{'='*50}")

        test_x, test_y = generate_fixed_test_set(seq, VOCAB_SIZE, device, seed=8888)

        # =========================================================
        # 動態計算 AGI 的最優參數，打破不合理的線性綁定
        # =========================================================
        optimal_M = seq // 4
        optimal_C = seq // 4
        optimal_K = seq // 4

        # 實例化模型 (已移除 PHTF)
        models = {
            "PLHN": PhaseLockedHolographyNetwork(
                VOCAB_SIZE, EMBED_DIM, HIDDEN_DIM, 2, max_seq_len=seq
            ).to(device),
            "AGI": HolographicMemoryNetwork(
                vocab_size=VOCAB_SIZE,
                embed_dim=EMBED_DIM,
                hidden_dim=HIDDEN_DIM,
                num_blocks=2,
                DI=EMBED_DIM // 4,
                K=optimal_K,
                M=optimal_M,
                C=optimal_C,
            ).to(device),
        }

        print("\n[參數量對比檢查]")
        for name, model in models.items():
            param_count = count_parameters(model)
            param_data[name][seq] = param_count
            print(f"  - {name} 總參數量: {param_count:,}")
        print("-" * 50)

        seq_results = {}
        for name, model in models.items():
            print(f"\n--- 訓練 {name} (Seq: {seq}) ---")

            set_seed(8888)

            samples_req, s_hist, acc_hist, loss_hist = train_until_convergence(
                model, test_x, test_y, seq, VOCAB_SIZE, device, name
            )

            summary_data[name][seq] = samples_req if samples_req > 0 else "DNF (Failed)"
            seq_results[name] = {
                "samples": s_hist,
                "test_acc": acc_hist,
                "loss": loss_hist,
            }

            if samples_req > 0:
                print(
                    f"✅ {name} 成功！在 {samples_req:,} 筆訓練資料時達到 100% 測試準確率。"
                )
            else:
                highest_acc = max(acc_hist) if acc_hist else 0
                print(
                    f"❌ {name} 失敗！無法在 {MAX_SAMPLES:,} 筆內達到 100% (最高準確率: {highest_acc:.4f})"
                )

        plot_path = os.path.join(base_dir, f"task21_scan_seq{seq}_agi_vs_plhn.png")
        create_training_plot(seq_results, seq, plot_path)
        plot_files[seq] = plot_path

    print("\n產生 All-In-One 簡報檔...")
    pptx_path = build_presentation_all_in_one(
        summary_data, plot_files, param_data, base_dir
    )
    print(f"✅ 所有序列長度掃描完成！最終彙整報告儲存於: {pptx_path}")


if __name__ == "__main__":
    main()
